#!/usr/bin/env python3
"""
Fix the POTX template by adding body placeholders to layouts that are missing them.
This ensures all bullet layouts have proper body placeholders inherited from their master.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Inches
import shutil
import zipfile

def patch_potx(template_path):
    """Patch POTX to work with python-pptx."""
    temp_path = 'temp_fix_template.pptx'
    shutil.copy(template_path, temp_path)
    
    with zipfile.ZipFile(temp_path, 'r') as zip_read:
        content_types = zip_read.read('[Content_Types].xml').decode('utf-8')
        content_types = content_types.replace(
            'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
        )
        
        with zipfile.ZipFile(temp_path + '.new', 'w', zipfile.ZIP_DEFLATED) as zip_write:
            for item in zip_read.infolist():
                data = zip_read.read(item.filename)
                if item.filename == '[Content_Types].xml':
                    zip_write.writestr(item, content_types.encode('utf-8'))
                else:
                    zip_write.writestr(item, data)
    
    shutil.move(temp_path + '.new', temp_path)
    return temp_path

def add_body_placeholder_to_layout(layout, master):
    """Add a body placeholder to a layout if it doesn't have one."""
    # Check if layout already has a body placeholder
    has_body = False
    for shape in layout.placeholders:
        if shape.placeholder_format.type == 2:  # BODY type
            has_body = True
            break
    
    if has_body:
        print(f"  Layout '{layout.name}' already has body placeholder")
        return
    
    print(f"  Adding body placeholder to layout '{layout.name}'")
    
    # Get the body placeholder from the master as a reference
    master_body = None
    for shape in master.placeholders:
        if shape.placeholder_format.type == 2:  # BODY type
            master_body = shape
            break
    
    if not master_body:
        print(f"    ERROR: Master doesn't have a body placeholder!")
        return
    
    # Add body placeholder to the layout
    # Use similar position/size as master's body placeholder
    left = master_body.left
    top = master_body.top
    width = master_body.width
    height = master_body.height
    
    # Add placeholder to layout
    placeholder = layout.shapes.add_placeholder(2, left, top, width, height)  # 2 = BODY type
    print(f"    Added body placeholder: left={left}, top={top}, width={width}, height={height}")

def main():
    template_path = 'templates/4734_template.potx'
    output_path = 'templates/4734_template_fixed.potx'
    
    print(f"Fixing template: {template_path}")
    
    # Patch and load template
    temp_path = patch_potx(template_path)
    prs = Presentation(temp_path)
    
    print(f"\nFound {len(prs.slide_masters)} masters")
    
    # Fix each master's layouts
    layouts_fixed = 0
    for i, master in enumerate(prs.slide_masters):
        print(f"\nMaster {i}: {len(master.slide_layouts)} layouts")
        
        # Check if master has body placeholder
        has_master_body = any(s.placeholder_format.type == 2 for s in master.placeholders if s.is_placeholder)
        
        if not has_master_body:
            print(f"  Master {i} has no body placeholder - skipping")
            continue
        
        for layout in master.slide_layouts:
            # Check if this is a bullet layout that needs fixing
            if 'Bullet' in layout.name or 'Photo' in layout.name:
                add_body_placeholder_to_layout(layout, master)
                layouts_fixed += 1
    
    # Save the fixed template
    print(f"\nSaving fixed template to: {output_path}")
    prs.save(output_path)
    
    # Convert back to POTX format
    print("Converting back to POTX format...")
    with zipfile.ZipFile(output_path, 'r') as zip_read:
        content_types = zip_read.read('[Content_Types].xml').decode('utf-8')
        content_types = content_types.replace(
            'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml',
            'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml'
        )
        
        with zipfile.ZipFile(output_path + '.new', 'w', zipfile.ZIP_DEFLATED) as zip_write:
            for item in zip_read.infolist():
                data = zip_read.read(item.filename)
                if item.filename == '[Content_Types].xml':
                    zip_write.writestr(item, content_types.encode('utf-8'))
                else:
                    zip_write.writestr(item, data)
    
    shutil.move(output_path + '.new', output_path)
    
    print(f"\nDone! Fixed {layouts_fixed} layouts")
    print(f"Backup your original template, then replace it with: {output_path}")

if __name__ == '__main__':
    main()
