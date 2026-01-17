import json
import os
from pptx import Presentation
from pptx_builder import build_pptx_from_slides


def test_title_slide_does_not_include_topics_in_pptx(tmp_path):
    # Create a single title slide
    slides_data = [
        (None, '', None, None, None, None, None, True, False, False, False, 'title')
    ]

    output_file = str(tmp_path / 'title_topics_test.pptx')

    success = build_pptx_from_slides(
        slides_data=slides_data,
        output_path=output_file,
        template_path='templates/4734_template.potx',
        pptx_layouts_map=json.load(open('pptx_layouts.json')),
        deck_info={'course_title': 'Test Course', 'week': 'Week X', 'date': '2099-12-31'}
    )

    assert success
    assert os.path.exists(output_file)

    prs = Presentation(output_file)
    assert len(prs.slides) == 1

    slide = prs.slides[0]

    # Collect text from all text-bearing shapes
    all_text = []
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            try:
                txt = shape.text
                if txt and txt.strip():
                    all_text.append(txt.strip())
            except Exception:
                pass

    combined = '\n'.join(all_text)

    assert 'Topic Alpha' not in combined
    assert 'Topic Beta' not in combined
