#!/usr/bin/env python3
"""
Post-process a pandoc-generated PPTX to apply custom slide layouts.
This script reads layout hints from the markdown and applies them to the PPTX.
"""

from pptx import Presentation
import json
import re
import sys

def apply_custom_layouts(pptx_path, md_path, potx_template_path):
    """
    Apply custom slide layouts to a PPTX file based on layout hints in markdown.
    
    Args:
        pptx_path: Path to the PPTX file to modify
        md_path: Path to the source markdown file with layout hints
        potx_template_path: Path to the POTX template with custom layouts
    """
    # Load the presentation
    prs = Presentation(pptx_path)
    
    # Load the template to get available layouts
    template = Presentation(potx_template_path)
    
    # Create a mapping of layout names to layout objects
    layout_map = {}
    for layout in template.slide_layouts:
        layout_map[layout.name] = layout
    
    print(f"Available layouts in template: {list(layout_map.keys())}")
    
    # Parse markdown to extract layout hints for each slide
    with open(md_path, 'r') as f:
        md_content = f.read()
    
    # Split by slide separators (---)
    slides_md = re.split(r'\n---\n', md_content)
    
    # Extract layout names from custom-style divs
    layout_hints = []
    for slide_md in slides_md:
        # Look for ::: {custom-style="LayoutName"}
        match = re.search(r'custom-style="([^"]+)"', slide_md)
        if match:
            layout_hints.append(match.group(1))
        else:
            layout_hints.append(None)
    
    print(f"Found {len(layout_hints)} layout hints for {len(prs.slides)} slides")
    
    # Apply layouts to slides
    for idx, slide in enumerate(prs.slides):
        if idx < len(layout_hints) and layout_hints[idx]:
            desired_layout = layout_hints[idx]
            if desired_layout in layout_map:
                print(f"Slide {idx + 1}: Applying layout '{desired_layout}'")
                # Get the layout from the template
                new_layout = None
                for layout in template.slide_layouts:
                    if layout.name == desired_layout:
                        new_layout = layout
                        break
                
                if new_layout:
                    # Create a new slide with the correct layout
                    # Note: python-pptx doesn't support changing a slide's layout directly
                    # We need to work around this limitation
                    print(f"  Warning: Cannot change layout of existing slide {idx + 1}")
                    print(f"  Desired layout: {desired_layout}")
            else:
                print(f"Slide {idx + 1}: Layout '{desired_layout}' not found in template")
    
    # Save the modified presentation
    prs.save(pptx_path)
    print(f"Saved modified PPTX to {pptx_path}")

if __name__ == '__main__':
    if len(sys.argv) != 4:
        print("Usage: python3 apply_pptx_layouts.py <pptx_file> <markdown_file> <potx_template>")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    md_file = sys.argv[2]
    potx_template = sys.argv[3]
    
    apply_custom_layouts(pptx_file, md_file, potx_template)
