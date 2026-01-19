from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import zipfile
import tempfile
import os

def convert_potx_to_pptx(potx_path):
    """Convert POTX to PPTX by modifying content type"""
    with tempfile.NamedTemporaryFile(mode='w+b', suffix='.pptx', delete=False) as tmp:
        tmp_path = tmp.name
    
    with zipfile.ZipFile(potx_path, 'r') as zip_in:
        with zipfile.ZipFile(tmp_path, 'w') as zip_out:
            for item in zip_in.infolist():
                data = zip_in.read(item.filename)
                if item.filename == '[Content_Types].xml':
                    data = data.replace(
                        b'application/vnd.openxmlformats-officedocument.presentationml.template.main',
                        b'application/vnd.openxmlformats-officedocument.presentationml.presentation.main'
                    )
                zip_out.writestr(item, data)
    
    return tmp_path

# Load template
template_path = 'templates/4734_template.potx'
tmp_pptx = convert_potx_to_pptx(template_path)

try:
    prs = Presentation(tmp_pptx)
    
    # Check White_Bullets and Gold_Bullets layouts for ALL shapes
    layouts_to_check = ['White_Bullets', 'Gold_Bullets']
    
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name in layouts_to_check:
                print(f"\nLayout: {layout.name}")
                print(f"  Total shapes: {len(layout.shapes)}")
                
                for shape_idx, shape in enumerate(layout.shapes):
                    print(f"\n  Shape {shape_idx}:")
                    print(f"    Type: {shape.shape_type}")
                    try:
                        print(f"    Name: {shape.name}")
                    except:
                        print(f"    Name: CANNOT_DETERMINE")
                    
                    # Check if it's a placeholder
                    is_placeholder = shape.is_placeholder if hasattr(shape, 'is_placeholder') else False
                    print(f"    Is placeholder: {is_placeholder}")
                    
                    # Check if it has text frame
                    if hasattr(shape, 'text_frame'):
                        print(f"    Has text_frame: True")
                        print(f"    Text: '{shape.text}'")
                    else:
                        print(f"    Has text_frame: False")

finally:
    os.unlink(tmp_pptx)
