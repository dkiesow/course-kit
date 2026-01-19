#!/usr/bin/env python3
"""
Create a PowerPoint template (.potx) with custom master slides matching our CSS styling.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme from CSS
COLORS = {
    'gold': RGBColor(228, 180, 68),  # #E4B444
    'white': RGBColor(255, 255, 255),
    'black': RGBColor(0, 0, 0),
    'dark_gray': RGBColor(51, 51, 51),
}

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)

def create_text_box(slide, left, top, width, height, text="", font_size=44, bold=True, 
                    color=COLORS['black'], align=PP_ALIGN.LEFT):
    """Helper to create a text box with standard styling"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.alignment = align
    p.font.name = 'Arial'  # Fallback - PowerPoint will use installed fonts
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    
    return textbox

def set_background(slide, color):
    """Set solid background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_colorbar(slide, color=COLORS['gold']):
    """Add the colorbar at top of slide"""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(0.15)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def create_template_bullets_master(prs):
    """Create master for template-bullets layout"""
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # White background
    set_background(slide, COLORS['white'])
    
    # Colorbar at top
    add_colorbar(slide)
    
    # Title placeholder
    create_text_box(slide, Inches(0.5), Inches(0.8), Inches(9), Inches(0.8),
                   "Headline Text", font_size=44, bold=True, color=COLORS['black'])
    
    # Bullet points placeholder
    create_text_box(slide, Inches(0.5), Inches(2), Inches(9), Inches(4.5),
                   "• Bullet point one\n• Bullet point two\n• Bullet point three", 
                   font_size=32, bold=False, color=COLORS['dark_gray'])
    
    return slide

def create_gold_bullets_master(prs):
    """Create master for gold-bullets layout"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Gold background
    set_background(slide, COLORS['gold'])
    
    # Title placeholder
    create_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8),
                   "Headline Text", font_size=44, bold=True, color=COLORS['white'])
    
    # Bullet points placeholder
    create_text_box(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(4.5),
                   "• Bullet point one\n• Bullet point two\n• Bullet point three", 
                   font_size=32, bold=False, color=COLORS['white'])
    
    return slide

def create_bullets_image_master(prs):
    """Create master for template-bullets-image layout"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # White background
    set_background(slide, COLORS['white'])
    
    # Colorbar
    add_colorbar(slide)
    
    # Title
    create_text_box(slide, Inches(0.5), Inches(0.8), Inches(9), Inches(0.8),
                   "Headline Text", font_size=44, bold=True, color=COLORS['black'])
    
    # Bullets on left
    create_text_box(slide, Inches(0.5), Inches(2), Inches(4.5), Inches(4.5),
                   "• Bullet one\n• Bullet two\n• Bullet three", 
                   font_size=28, bold=False, color=COLORS['dark_gray'])
    
    # Image placeholder on right
    create_text_box(slide, Inches(5.5), Inches(2), Inches(4), Inches(4.5),
                   "[Image]", font_size=32, bold=False, color=COLORS['dark_gray'],
                   align=PP_ALIGN.CENTER)
    
    return slide

def create_photo_centered_master(prs):
    """Create master for photo-centered layout"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # White background
    set_background(slide, COLORS['white'])
    
    # Colorbar
    add_colorbar(slide)
    
    # Title at top
    create_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.6),
                   "Headline Text", font_size=36, bold=True, color=COLORS['black'])
    
    # Large centered image area
    create_text_box(slide, Inches(1.5), Inches(2), Inches(7), Inches(4.5),
                   "[Centered Image]", font_size=40, bold=False, color=COLORS['dark_gray'],
                   align=PP_ALIGN.CENTER)
    
    return slide

def create_quote_master(prs):
    """Create master for quote layouts"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # White background
    set_background(slide, COLORS['white'])
    
    # Colorbar
    add_colorbar(slide)
    
    # Title
    create_text_box(slide, Inches(0.5), Inches(0.8), Inches(9), Inches(0.8),
                   "Headline Text", font_size=44, bold=True, color=COLORS['black'])
    
    # Quote text
    create_text_box(slide, Inches(1), Inches(2.5), Inches(8), Inches(2.5),
                   '"Quote text here"', font_size=36, bold=False, color=COLORS['dark_gray'],
                   align=PP_ALIGN.LEFT)
    
    # Attribution
    create_text_box(slide, Inches(1), Inches(5.2), Inches(8), Inches(0.6),
                   '— Attribution', font_size=24, bold=False, color=COLORS['dark_gray'],
                   align=PP_ALIGN.RIGHT)
    
    return slide

def create_gold_quote_master(prs):
    """Create master for gold-quote layout"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Gold background
    set_background(slide, COLORS['gold'])
    
    # Title
    create_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8),
                   "Headline Text", font_size=44, bold=True, color=COLORS['white'])
    
    # Quote text
    create_text_box(slide, Inches(1), Inches(2), Inches(8), Inches(2.5),
                   '"Quote text here"', font_size=36, bold=False, color=COLORS['white'],
                   align=PP_ALIGN.LEFT)
    
    # Attribution
    create_text_box(slide, Inches(1), Inches(5), Inches(8), Inches(0.6),
                   '— Attribution', font_size=24, bold=False, color=COLORS['white'],
                   align=PP_ALIGN.RIGHT)
    
    return slide

def create_closing_master(prs):
    """Create master for closing slide"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Gold background
    set_background(slide, COLORS['gold'])
    
    # Contact info (multi-line headline)
    create_text_box(slide, Inches(1), Inches(1.5), Inches(8), Inches(3),
                   "Contact Name\nTitle\nLocation\nemail@example.com", 
                   font_size=32, bold=True, color=COLORS['white'],
                   align=PP_ALIGN.CENTER)
    
    # Thank you message
    create_text_box(slide, Inches(1), Inches(4.5), Inches(8), Inches(1),
                   "Thank You", font_size=48, bold=True, color=COLORS['white'],
                   align=PP_ALIGN.CENTER)
    
    # Add note about logo placement
    create_text_box(slide, Inches(3), Inches(6.2), Inches(4), Inches(0.5),
                   "[Logo Here]", font_size=20, bold=False, color=COLORS['white'],
                   align=PP_ALIGN.CENTER)
    
    return slide

def create_title_master(prs):
    """Create master for title slide"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add background image placeholder (white for now)
    set_background(slide, RGBColor(240, 240, 240))
    
    # Colorbar
    add_colorbar(slide)
    
    # Course title
    create_text_box(slide, Inches(1), Inches(2.5), Inches(8), Inches(1),
                   "Course Title", font_size=52, bold=True, color=COLORS['black'],
                   align=PP_ALIGN.LEFT)
    
    # Week/Topic
    create_text_box(slide, Inches(1), Inches(3.7), Inches(8), Inches(0.8),
                   "Week Topic", font_size=40, bold=True, color=COLORS['dark_gray'],
                   align=PP_ALIGN.LEFT)
    
    # Date
    create_text_box(slide, Inches(1), Inches(4.7), Inches(8), Inches(0.6),
                   "Date", font_size=32, bold=False, color=COLORS['dark_gray'],
                   align=PP_ALIGN.LEFT)
    
    return slide

def create_lines_master(prs):
    """Create master for template-lines layout"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # White background
    set_background(slide, COLORS['white'])
    
    # Colorbar
    add_colorbar(slide)
    
    # Title
    create_text_box(slide, Inches(0.5), Inches(0.8), Inches(9), Inches(0.8),
                   "Headline Text", font_size=44, bold=True, color=COLORS['black'])
    
    # Text content (no bullets, just lines)
    create_text_box(slide, Inches(0.5), Inches(2), Inches(9), Inches(4.5),
                   "Line one\nLine two\nLine three", 
                   font_size=32, bold=False, color=COLORS['dark_gray'])
    
    return slide

def main():
    """Create the master template POTX file"""
    print("Creating PowerPoint master template...")
    
    prs = Presentation()
    
    # Create example slides for each template type
    print("Creating Title slide master...")
    create_title_master(prs)
    
    print("Creating template-bullets master...")
    create_template_bullets_master(prs)
    
    print("Creating gold-bullets master...")
    create_gold_bullets_master(prs)
    
    print("Creating template-bullets-image master...")
    create_bullets_image_master(prs)
    
    print("Creating photo-centered master...")
    create_photo_centered_master(prs)
    
    print("Creating template-quote master...")
    create_quote_master(prs)
    
    print("Creating gold-quote master...")
    create_gold_quote_master(prs)
    
    print("Creating template-lines master...")
    create_lines_master(prs)
    
    print("Creating closing slide master...")
    create_closing_master(prs)
    
    # Save as PPTX (can be saved as POTX manually in PowerPoint)
    output_path = 'output/classroom_master_slides.pptx'
    os.makedirs('output', exist_ok=True)
    prs.save(output_path)
    
    print(f"\n✅ Success! Created {output_path}")
    print("\nThis template includes master slides for:")
    print("  1. Title slide")
    print("  2. Template-bullets (white with colorbar)")
    print("  3. Gold-bullets (gold background)")
    print("  4. Template-bullets-image (split layout)")
    print("  5. Photo-centered")
    print("  6. Template-quote")
    print("  7. Gold-quote")
    print("  8. Template-lines")
    print("  9. Closing slide")
    print("\nYou can now:")
    print("1. Open this PPTX in PowerPoint")
    print("2. View → Slide Master to customize layouts")
    print("3. Add background images and adjust fonts/colors")
    print("4. File → Save As → PowerPoint Template (.potx) to save as template")

if __name__ == '__main__':
    main()
