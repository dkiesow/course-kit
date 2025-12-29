from pptx import Presentation
import zipfile, shutil, os

# Copy and convert POTX to PPTX
shutil.copy2('4734_template.potx', 'temp_check.pptx')
z = zipfile.ZipFile('temp_check.pptx', 'r')
files = {n: z.read(n) for n in z.namelist()}
z.close()

content_xml = files['[Content_Types].xml'].decode('utf-8')
content_xml = content_xml.replace('presentationml.template.main', 'presentationml.presentation.main')
files['[Content_Types].xml'] = content_xml.encode('utf-8')

z = zipfile.ZipFile('temp_check.pptx', 'w')
for n, d in files.items():
    z.writestr(n, d)
z.close()

prs = Presentation('temp_check.pptx')

# Check Master 0 (White theme)
master = prs.slide_masters[0]
print("=== Master 0 (White) Shapes ===")
for i, shape in enumerate(master.shapes):
    print(f"Shape {i}:")
    print(f"  Type: {shape.shape_type}")
    print(f"  Has text frame: {shape.has_text_frame}")
    if shape.is_placeholder:
        print(f"  Is placeholder: True")
        print(f"  Placeholder idx: {shape.placeholder_format.idx}")
        print(f"  Placeholder type: {shape.placeholder_format.type}")
    if hasattr(shape, 'name'):
        print(f"  Name: {shape.name}")
    print()

# Check White_Bullets layout
print("\n=== White_Bullets Layout Shapes ===")
white_bullets_layout = prs.slide_masters[0].slide_layouts[0]
for i, shape in enumerate(white_bullets_layout.shapes):
    print(f"Shape {i}:")
    print(f"  Type: {shape.shape_type}")
    print(f"  Has text frame: {shape.has_text_frame}")
    if shape.is_placeholder:
        print(f"  Is placeholder: True")
        print(f"  Placeholder idx: {shape.placeholder_format.idx}")
        print(f"  Placeholder type: {shape.placeholder_format.type}")
    if hasattr(shape, 'name'):
        print(f"  Name: {shape.name}")
    print()

# Create an actual slide and check what's available
print("\n=== Slide created from White_Bullets ===")
slide = prs.slides.add_slide(white_bullets_layout)
for i, shape in enumerate(slide.shapes):
    print(f"Shape {i}:")
    print(f"  Type: {shape.shape_type}")
    print(f"  Has text frame: {shape.has_text_frame}")
    if shape.is_placeholder:
        print(f"  Is placeholder: True")
        print(f"  Placeholder idx: {shape.placeholder_format.idx}")
        print(f"  Placeholder type: {shape.placeholder_format.type}")
    if hasattr(shape, 'name'):
        print(f"  Name: {shape.name}")
    print()

# Check placeholders collection
print("\n=== Slide.placeholders collection ===")
for idx in slide.placeholders:
    placeholder = slide.placeholders[idx]
    print(f"Placeholder idx {idx}:")
    print(f"  Type: {placeholder.placeholder_format.type}")
    print(f"  Has text frame: {placeholder.has_text_frame}")
    if hasattr(placeholder, 'name'):
        print(f"  Name: {placeholder.name}")
    print()

os.remove('temp_check.pptx')
