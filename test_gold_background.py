#!/usr/bin/env python3
"""
Test if slides created with Gold layouts show the gold background.
"""

from pptx import Presentation
import shutil
import zipfile

# Patch POTX
template_path = '4734_template.potx'
temp_path = 'temp_template.pptx'

shutil.copy(template_path, temp_path)

with zipfile.ZipFile(temp_path, 'r') as zip_read:
    content_types = zip_read.read('[Content_Types].xml').decode('utf-8')
    content_types = content_types.replace(
        'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
    )
    
    with zipfile.ZipFile(temp_path + '.new', 'w', zipfile.ZIP_DEFLATED) as zip_write:
        for item in zip_read.infolist():
            data = zip_read.read(item.filename)
            if item.filename == '[Content_Types].xml':
                zip_write.writestr(item, content_types.encode('utf-8'))
            else:
                zip_write.writestr(item, data)

shutil.move(temp_path + '.new', temp_path)

prs = Presentation(temp_path)

# Get all layouts
layouts = {}
for master in prs.slide_masters:
    for layout in master.slide_layouts:
        layouts[layout.name] = layout

print("Creating test slides...")

# Create slide with White_Bullets
if 'White_Bullets' in layouts:
    slide = prs.slides.add_slide(layouts['White_Bullets'])
    title = slide.shapes.title
    title.text = "White Bullets Test"
    print("  Created White_Bullets slide")

# Create slide with Gold_Bullets
if 'Gold_Bullets' in layouts:
    slide = prs.slides.add_slide(layouts['Gold_Bullets'])
    title = slide.shapes.title
    title.text = "Gold Bullets Test"
    print("  Created Gold_Bullets slide")

# Save
output_path = 'test_gold_background.pptx'
prs.save(output_path)
print(f"\nSaved: {output_path}")
print("Open it to check if Gold slide has gold background")
