#!/usr/bin/env python3
"""Inspect placeholder structures in all layouts."""

import zipfile
import shutil
from pptx import Presentation

# Copy and patch template
shutil.copy('4734_template.potx', 'temp_check.pptx')

try:
    # Patch content type
    with zipfile.ZipFile('temp_check.pptx', 'r') as zf:
        files = {name: zf.read(name) for name in zf.namelist()}
    
    content_xml = files['[Content_Types].xml'].decode('utf-8')
    content_xml = content_xml.replace(
        'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
    )
    files['[Content_Types].xml'] = content_xml.encode('utf-8')
    
    with zipfile.ZipFile('temp_check.pptx', 'w', zipfile.ZIP_DEFLATED) as zf:
        for name, data in files.items():
            zf.writestr(name, data)
except Exception as e:
    print(f"Warning: Could not patch template: {e}")

# Load and inspect placeholders
prs = Presentation('temp_check.pptx')

print("Placeholder Structure Analysis\n")
print("=" * 80)

for master_idx, master in enumerate(prs.slide_masters):
    master_name = master.name if hasattr(master, 'name') else f'Master {master_idx}'
    print(f"\n{master_name}")
    print("-" * 80)
    
    for layout_idx, layout in enumerate(master.slide_layouts):
        print(f"\n  Layout: {layout.name}")
        
        if len(layout.placeholders) == 0:
            print(f"    No placeholders")
        else:
            for ph in layout.placeholders:
                # Get the actual type value
                ph_type = ph.placeholder_format.type
                
                # Standard placeholder type names (from MSO_PLACEHOLDER_TYPE enum)
                placeholder_type_name = {
                    1: 'TITLE',
                    2: 'BODY', 
                    3: 'CENTERED_TITLE',
                    4: 'SUBTITLE',
                    5: 'DATE',
                    6: 'SLIDE_NUMBER',
                    7: 'FOOTER',
                    8: 'HEADER',
                    9: 'OBJECT',
                    10: 'CHART',
                    11: 'TABLE',
                    12: 'CLIP_ART',
                    13: 'DIAGRAM',
                    14: 'MEDIA',
                    15: 'SLIDE_IMAGE',
                    16: 'PICTURE',
                    17: 'CONTENT',
                    18: 'PICTURE'
                }.get(ph_type, f'UNKNOWN_{ph_type}')
                
                has_text_frame = hasattr(ph, 'text_frame') and ph.text_frame is not None
                
                # Get shape type info
                shape_type = ph.shape_type if hasattr(ph, 'shape_type') else 'N/A'
                
                print(f"    IDX:{ph.placeholder_format.idx} Type:{placeholder_type_name}({ph_type}) TextFrame:{has_text_frame} Shape:{shape_type}")

print("\n" + "=" * 80)
