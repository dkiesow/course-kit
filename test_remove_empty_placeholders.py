import json
import os
from pptx import Presentation
from pptx_builder import build_pptx_from_slides


def test_remove_empty_body_placeholders(tmp_path):
    # One image-only slide (no headline/paragraph/bullets) using an image layout
    slides_data = [
        (None, '', None, None, None, None, 'milkshake.png', 0, False, False, False, 'template-bullets-image')
    ]

    output_file = str(tmp_path / "test_out.pptx")

    success = build_pptx_from_slides(
        slides_data=slides_data,
        output_path=output_file,
        template_path='templates/4734_template.potx',
        pptx_layouts_map=json.load(open('pptx_layouts.json'))
    )

    assert success
    assert os.path.exists(output_file)

    prs = Presentation(output_file)
    assert len(prs.slides) == 1

    slide = prs.slides[0]

    # Ensure there are no empty Body/Object/Content placeholders
    for shape in slide.shapes:
        if getattr(shape, 'is_placeholder', False) and shape.has_text_frame:
            try:
                ph_type = shape.placeholder_format.type
            except Exception:
                continue
            if ph_type in [2, 7, 14]:
                # The placeholder should contain non-empty text if present
                assert shape.text.strip() != '', f"Found empty content placeholder: {getattr(shape, 'name', None)}"
