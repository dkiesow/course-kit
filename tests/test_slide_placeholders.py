from pptx import Presentation
import zipfile, shutil, os

# Copy and convert POTX to PPTX
shutil.copy2('templates/4734_template.potx', 'temp_check.pptx')
z = zipfile.ZipFile('temp_check.pptx', 'r')
files = {n: z.read(n) for n in z.namelist()}
z.close()

# Fix content type
content_xml = files['[Content_Types].xml'].decode('utf-8')
content_xml = content_xml.replace('presentationml.template.main', 'presentationml.presentation.main')
files['[Content_Types].xml'] = content_xml.encode('utf-8')

z = zipfile.ZipFile('temp_check.pptx', 'w')
for n, d in files.items():
    z.writestr(n, d)
z.close()

# Load presentation
prs = Presentation('temp_check.pptx')

# Test by actually creating slides and checking their placeholders
print("Testing actual slide placeholders (inherited from master + layout):\n")

# Test a few key layouts
test_layouts = ['White_Bullets', 'Gold_Bullets', 'Arches_Title', 'White_Bullets_Photo']

for layout_name in test_layouts:
    # Find the layout
    layout = None
    for master in prs.slide_masters:
        for l in master.slide_layouts:
            if l.name == layout_name:
                layout = l
                break
        if layout:
            break
    
    if not layout:
        print(f"Layout '{layout_name}' not found")
        continue
    
    # Create an actual slide
    slide = prs.slides.add_slide(layout)
    
    print(f"Layout: {layout_name}")
    print(f"  Slide has {len(slide.placeholders)} placeholders:")
    for ph in slide.placeholders:
        ph_type = ph.placeholder_format.type
        print(f"    idx={ph.placeholder_format.idx}, type={ph_type}, has_text_frame={ph.has_text_frame}")
    print()

# Cleanup
os.remove('temp_check.pptx')
