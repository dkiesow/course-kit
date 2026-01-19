import json
import os
from pptx import Presentation
from pptx_builder import build_pptx_from_slides


def test_larger_image_constrained_to_placeholder(tmp_path):
    image_name = 'milkshake.png'
    assert os.path.exists(os.path.join('assets', image_name)), "Test image not found in assets"

    slides_large = [
        (None, 'Headline', None, None, None, None, image_name, 0, False, True, False, 'bullets-image-top')
    ]
    out_large = str(tmp_path / 'large_placeholder.pptx')
    build_pptx_from_slides(slides_large, out_large, 'templates/4734_template.potx', json.load(open('pptx_layouts.json')))

    prs = Presentation(out_large)
    assert len(prs.slides) == 1

    slide = prs.slides[0]

    # Find the first picture
    pic = None
    for shape in slide.shapes:
        if shape.shape_type == 13:  # picture
            pic = shape
            break
    assert pic is not None, "No picture found on slide"

    # Find picture placeholder geometry from the slide's layout
    ph = None
    try:
        for placeholder in slide.slide_layout.placeholders:
            if placeholder.placeholder_format.type == 18:  # picture placeholder
                ph = placeholder
                break
    except Exception:
        # If layout doesn't expose placeholders, fail the test
        ph = None

    assert ph is not None, "No picture placeholder found in layout"

    # Ensure picture fits inside placeholder bounds
    assert pic.width <= ph.width, f"Picture width {pic.width} exceeds placeholder {ph.width}"
    assert pic.left >= ph.left, f"Picture left {pic.left} is left of placeholder {ph.left}"
    assert pic.left + pic.width <= ph.left + ph.width, f"Picture right exceeds placeholder right"

    # Ensure it's horizontally centered within reasonable tolerance (2000 EMU ~ 0.002 in)
    expected_left = ph.left + (ph.width - pic.width) // 2
    assert abs(pic.left - expected_left) <= 2000, f"Picture not centered (expected {expected_left}, got {pic.left})"