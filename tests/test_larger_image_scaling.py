import json
import os
from pptx import Presentation
from pptx_builder import build_pptx_from_slides


def get_first_picture_size(pptx_path):
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # picture
                return shape.width, shape.height
    return None, None


def test_larger_image_increases_picture_size(tmp_path):
    image_name = 'milkshake.png'
    assert os.path.exists(os.path.join('assets', image_name)), "Test image not found in assets"

    slides_small = [
        (None, 'Headline', None, None, None, None, image_name, 0, False, False, False, 'bullets-image-top')
    ]
    out_small = str(tmp_path / 'small.pptx')
    build_pptx_from_slides(slides_small, out_small, 'templates/4734_template.potx', json.load(open('pptx_layouts.json')))
    w_small, h_small = get_first_picture_size(out_small)

    slides_large = [
        (None, 'Headline', None, None, None, None, image_name, 0, False, True, False, 'bullets-image-top')
    ]
    out_large = str(tmp_path / 'large.pptx')
    build_pptx_from_slides(slides_large, out_large, 'templates/4734_template.potx', json.load(open('pptx_layouts.json')))
    w_large, h_large = get_first_picture_size(out_large)

    assert w_small is not None and w_large is not None, "No pictures found in generated PPTX"
    # Expect the larger-image variant to have a width greater than the small variant
    assert w_large > w_small, f"Expected larger image width (got {w_large} <= {w_small})"