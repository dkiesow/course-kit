#!/usr/bin/env python3
"""
Custom PPTX builder using python-pptx to work with custom layout names.
This replaces pandoc for PPTX generation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import json
import os
import sys
import tempfile
import shutil
import zipfile
import xml.etree.ElementTree as ET

# Threshold (in bytes) above which animated GIFs are converted to PNG fallbacks
GIF_FALLBACK_THRESHOLD = 5 * 1024 * 1024  # 5 MB


def build_pptx_from_slides(slides_data, output_path, template_path, pptx_layouts_map, deck_info=None):
    """
    Build a PPTX file directly from slide data using custom layouts.
    
    Args:
        slides_data: List of tuples containing slide data from database
        output_path: Path where PPTX should be saved
        template_path: Path to POTX/PPTX template file  
        pptx_layouts_map: Dict mapping template_base to layout names
        deck_info: Dict with course_title, week, date for title slide
    """
    import zipfile
    import shutil
    
    # python-pptx doesn't support POTX files, so we need to convert it
    temp_pptx = 'temp_template.pptx'
    
    # Remove old temp file if it exists to ensure we use latest template
    if os.path.exists(temp_pptx):
        os.remove(temp_pptx)
    
    # Copy the template
    shutil.copy2(template_path, temp_pptx)
    
    # Modify the content type in [Content_Types].xml to make it a presentation
    try:
        # Read the ZIP
        with zipfile.ZipFile(temp_pptx, 'r') as zf:
            files = {name: zf.read(name) for name in zf.namelist()}
        
        # Modify content types
        if '[Content_Types].xml' in files:
            content_xml = files['[Content_Types].xml'].decode('utf-8')
            content_xml = content_xml.replace(
                'presentationml.template.main',
                'presentationml.presentation.main'
            )
            # Add GIF support if not already present
            if 'image/gif' not in content_xml:
                # Insert GIF content type after PNG
                content_xml = content_xml.replace(
                    '<Default Extension="png" ContentType="image/png"/>',
                    '<Default Extension="png" ContentType="image/png"/><Default Extension="gif" ContentType="image/gif"/>'
                )
                print("Added GIF content type support to template")
            files['[Content_Types].xml'] = content_xml.encode('utf-8')
        
        # Write back
        with zipfile.ZipFile(temp_pptx, 'w', zipfile.ZIP_DEFLATED) as zf:
            for name, data in files.items():
                zf.writestr(name, data)
    except Exception as e:
        print(f"Warning: Could not patch template: {e}")
    
    # Now load with python-pptx
    prs = Presentation(temp_pptx)
    
    # Build a map of layout names to layout objects
    # Must iterate through all slide_masters since prs.slide_layouts only returns first master's layouts
    layout_map = {}
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            layout_map[layout.name] = layout
    
    print(f"Available layouts: {list(layout_map.keys())}")
    
    # Process each slide
    for slide_data in slides_data:
        (slide_class, headline, paragraph, bullets, quote, quote_citation, 
         image_path, is_title, hide_headline, larger_image, fullscreen, template_base) = slide_data
        
        # Determine layout name directly from slide_class (no guessing)
        if is_title:
            layout_name = "Arches_Title"
        else:
            # Use slide_class if it exists, otherwise template_base
            template_key = slide_class if slide_class else template_base
            
            import sys
            print(f"  Slide template_key: {template_key}, quote: {repr(quote[:50] if quote else None)}, image_path: {repr(image_path)}", file=sys.stderr, flush=True)
            
            # For bullets-image-top templates with fullscreen, switch to full-photo-headline
            if template_key == 'bullets-image-top' and fullscreen:
                template_key = 'full-photo-headline'
            elif template_key == 'gold-bullets-image-top' and fullscreen:
                template_key = 'gold-full-photo-headline'
            
            layout_name = pptx_layouts_map.get(template_key)
            
            # For photo-centered templates, use headline version if hide_headline is False
            if template_key == 'photo-centered' and not hide_headline:
                layout_name = "White_Photo_Headline"
            elif template_key == 'gold-photo-centered' and not hide_headline:
                layout_name = "Gold_Photo_Headline"
            # For bullets-image-top templates, use big photo version if larger_image is True
            elif template_key == 'bullets-image-top' and larger_image:
                layout_name = "White_Top_Bullets_Big_Photo"
            elif template_key == 'gold-bullets-image-top' and larger_image:
                layout_name = "Gold_Top_Bullets_Big_Photo"
            
            if not layout_name:
                print(f"Warning: No layout mapping for template '{template_key}'")
                layout_name = "White_Bullets"
        
        # Get the layout
        layout = layout_map.get(layout_name)
        if not layout:
            print(f"Warning: Layout '{layout_name}' not found, using first available layout")
            layout = prs.slide_layouts[0]
        
        # Add slide with the layout
        slide = prs.slides.add_slide(layout)
        
        # Populate placeholders based on slide type
        if is_title:
            populate_title_slide(slide, headline, paragraph, deck_info)
        elif template_key == "closing":
            populate_closing_slide(slide, headline, paragraph, image_path)
        elif template_key in ["photo-centered", "gold-photo-centered", "full-photo-headline", "gold-full-photo-headline", 
                              "template-photo-centered", "template-gold-photo-centered", "template-full-photo-headline", "template-gold-full-photo-headline"]:
            populate_photo_slide(slide, image_path, headline, hide_headline, paragraph, bullets, larger_image, prs)
        elif template_key in ["quote", "gold-quote", "template-quote", "template-gold-quote"]:
            populate_quote_slide(slide, quote, quote_citation)
        elif template_key in ["bullets-image", "bullets-image-split", "bullets-image-top", "gold-bullets-image-split", "gold-bullets-image-top",
                              "template-bullets-image", "template-bullets-image-split", "template-bullets-image-top", "template-gold-bullets-image-split", "template-gold-bullets-image-top"]:
            # Image layouts - add image
            populate_content_slide(slide, headline, paragraph, bullets, image_path, hide_headline, slide_class, larger_image, prs)
        else:
            # Regular content - no image unless it's explicitly an image layout
            populate_content_slide(slide, headline, paragraph, bullets, None, hide_headline, slide_class, larger_image, prs)
    
    # Save the presentation
    prs.save(output_path)
    # Normalize PPTX docProps (populate Slides/Words/etc) to keep document metadata accurate
    try:
        normalize_pptx(output_path)
    except Exception:
        pass
    return True


def add_formatted_text_to_frame(text_frame, text):
    """Add text with markdown formatting to a text frame."""
    import re
    
    # Clear existing paragraphs but keep the first one for formatting
    if text_frame.paragraphs:
        text_frame.paragraphs[0].text = ""
        while len(text_frame.paragraphs) > 1:
            elem = text_frame.paragraphs[1]._element
            elem.getparent().remove(elem)
    
    # Split text into lines and process each line
    lines = text.split('\n')
    
    for line_idx, line in enumerate(lines):
        if line.strip():  # Only process non-empty lines
            # Get or create paragraph
            if line_idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Parse markdown formatting in the line
            parse_markdown_to_paragraph(p, line)
        elif line_idx > 0:  # Add empty line breaks (but not at the beginning)
            text_frame.add_paragraph()


def parse_markdown_to_paragraph(paragraph, text):
    """Parse markdown formatting and add runs to paragraph."""
    import re
    
    # Pattern to match **bold**, *italic*, and `code` markdown
    # Match in order: code, bold, italic, or regular text
    # Use DOTALL flag so . matches newlines
    pattern = r'(`[^`]+`|\*\*[\s\S]*?\*\*|\*[\s\S]*?\*|[^*`]+|[*`])'
    parts = re.findall(pattern, text, re.DOTALL)
    
    for part in parts:
        if not part:
            continue
            
        if part.startswith('`') and part.endswith('`') and len(part) > 2:
            # Inline code
            run = paragraph.add_run()
            run.text = part[1:-1]  # Remove ` markers
            run.font.name = 'Courier New'
        elif part.startswith('**') and part.endswith('**') and len(part) > 4:
            # Bold text
            run = paragraph.add_run()
            run.text = part[2:-2]  # Remove ** markers
            run.font.bold = True
        elif part.startswith('*') and part.endswith('*') and len(part) > 2 and not part.startswith('**'):
            # Italic text
            run = paragraph.add_run()
            run.text = part[1:-1]  # Remove * markers
            run.font.italic = True
        else:
            # Regular text
            run = paragraph.add_run()
            run.text = part


def remove_empty_body_placeholders(slide):
    """Remove empty content/body placeholders from a slide.
    Returns a list of (placeholder_type, name) removed for debugging."""
    removed = []
    # Iterate a copy because we'll be removing shapes from the tree
    for shape in list(slide.shapes):
        if not getattr(shape, 'is_placeholder', False):
            continue
        try:
            ph_type = shape.placeholder_format.type
        except Exception:
            continue
        # Consider Body, Object, and Content placeholders
        if ph_type in [2, 7, 14] and getattr(shape, 'has_text_frame', False):
            txt = shape.text or ''
            if txt.strip() == '':
                try:
                    slide.shapes._spTree.remove(shape._element)
                    removed.append((ph_type, getattr(shape, 'name', None)))
                except Exception:
                    pass
    return removed


# --- PPTX normalization helpers (used by tests) ---

def _reorder_content_types(ct_bytes, files=None):
    """Reorder overrides in a [Content_Types].xml blob so that docProps entries occur
    before the presentation part. Returns a bytes object with the updated XML.
    """
    try:
        root = ET.fromstring(ct_bytes.decode('utf-8'))
    except Exception:
        return ct_bytes
    ns = root.tag.split('}')[0].strip('{')
    overrides = list(root.findall(f'{{{ns}}}Override'))

    others = []
    docprops = []
    presentation = []

    for o in overrides:
        pn = o.attrib.get('PartName')
        if pn and pn.startswith('/docProps/'):
            docprops.append(o)
        elif pn == '/ppt/presentation.xml':
            presentation.append(o)
        else:
            others.append(o)

    # Remove existing override elements
    for o in overrides:
        root.remove(o)

    # Append in the order: others, docprops, presentation
    for o in others + docprops + presentation:
        root.append(o)

    return ET.tostring(root, encoding='utf-8', xml_declaration=True)


def _normalize_slide_paragraph_pPr(files):
    """Ensure every paragraph (<a:p>) in slide XML has a <a:pPr> child.
    Modifies and returns a new files dict with updated slide XML where needed.
    """
    new_files = dict(files)
    for name, data in list(files.items()):
        if not name.startswith('ppt/slides/') or not name.endswith('.xml'):
            continue
        try:
            root = ET.fromstring(data)
        except Exception:
            continue

        changed = False
        # Find txBody elements and ensure each <a:p> contains <a:pPr>
        for el in root.iter():
            if el.tag.split('}', 1)[-1] == 'txBody':
                # iterate direct children paragraphs under txBody
                for p in el.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}p'):
                    # check for pPr child
                    if not any(ch.tag.split('}', 1)[-1] == 'pPr' for ch in list(p)):
                        pPr = ET.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
                        p.insert(0, pPr)
                        changed = True
        if changed:
            new_files[name] = ET.tostring(root, encoding='utf-8', xml_declaration=True)
    return new_files


def _compute_layout_master_overlap_subtract(files, used_layout_names, used_master_files):
    """Compute a small subtraction factor based on textual overlap between layouts and masters.
    Simple heuristic used in tests: count occurrences of a token like 'alpha' and compute min(layout_count, master_count) // 4.
    """
    layout_count = 0
    for path, data in files.items():
        if path.startswith('ppt/slideLayouts/'):
            s = data.decode('utf-8', errors='ignore')
            for lname in used_layout_names:
                if f'name="{lname}"' in s:
                    layout_count += s.count('alpha')
    master_count = 0
    for m in used_master_files:
        if m in files:
            master_count += files[m].decode('utf-8', errors='ignore').count('alpha')

    overlap = min(layout_count, master_count)
    return overlap // 4


def normalize_pptx(pptx_path):
    """Normalize a PPTX file by populating docProps (Words, Paragraphs, Slides, TitlesOfParts).
    This mutates the PPTX in place.
    """
    try:
        prs = Presentation(pptx_path)
    except Exception:
        return

    # Compute counts
    slide_count = len(prs.slides)
    words = 0
    paragraphs = 0
    titles = []
    for slide in prs.slides:
        # collect title if present
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                t = slide.shapes.title.text.strip()
                if t:
                    titles.append(t)
        except Exception:
            pass
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False):
                try:
                    text = shape.text or ''
                    words += len(text.split())
                    paragraphs += len(shape.text_frame.paragraphs)
                except Exception:
                    pass

    # Build app.xml content
    ns = 'http://schemas.openxmlformats.org/officeDocument/2006/extended-properties'
    ET.register_namespace('', ns)
    props = ET.Element(f'{{{ns}}}Properties')
    app = ET.SubElement(props, f'{{{ns}}}Application')
    app.text = 'python-pptx'
    w = ET.SubElement(props, f'{{{ns}}}Words')
    w.text = str(words)
    p = ET.SubElement(props, f'{{{ns}}}Paragraphs')
    p.text = str(paragraphs)
    s = ET.SubElement(props, f'{{{ns}}}Slides')
    s.text = str(slide_count)
    tops = ET.SubElement(props, f'{{{ns}}}TitlesOfParts')
    for t in titles:
        te = ET.SubElement(tops, f'{{{ns}}}t')
        te.text = t

    new_app_xml = ET.tostring(props, encoding='utf-8', xml_declaration=True)

    # Read and rewrite ZIP with updated docProps/app.xml
    tmpfd, tmpname = tempfile.mkstemp(suffix='.pptx')
    os.close(tmpfd)
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zin, zipfile.ZipFile(tmpname, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == 'docProps/app.xml':
                    zout.writestr(item, new_app_xml)
                else:
                    zout.writestr(item, data)
        shutil.move(tmpname, pptx_path)
    except Exception:
        try:
            os.remove(tmpname)
        except Exception:
            pass
    return


def populate_title_slide(slide, headline, paragraph, deck_info=None):
    """Populate title slide with course name, week, and date."""
    # Get deck information
    course_title = deck_info.get('course_title', 'Journalism Innovation') if deck_info else 'Journalism Innovation'
    week = deck_info.get('week', '') if deck_info else ''
    date = deck_info.get('date', '') if deck_info else ''
    
    # Find the three placeholders (should be course title, week, date)
    placeholders = [shape for shape in slide.placeholders if shape.placeholder_format.type in [1, 2, 7]]  # Title, Body, Object
    
    if len(placeholders) >= 3:
        # Placeholder 1: Course Title
        if placeholders[0].has_text_frame:
            placeholders[0].text = course_title
        
        # Placeholder 2: Week
        if placeholders[1].has_text_frame:
            placeholders[1].text = week if week else 'Week'
            
        # Placeholder 3: Date  
        if placeholders[2].has_text_frame:
            placeholders[2].text = date if date else 'Date'
    else:
        # Fallback to old behavior if not enough placeholders
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            if ph_type == 1:  # Title
                if shape.has_text_frame:
                    shape.text = course_title
                break


def populate_quote_slide(slide, quote, quote_citation):
    """Populate quote slide with quote text and citation."""
    import sys
    print(f"  populate_quote_slide - quote: {repr(quote)}, citation: {repr(quote_citation)}", file=sys.stderr, flush=True)
    
    # Find text placeholders - skip title placeholder (idx 0), use content placeholder
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Skip title placeholders
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 1:  # 1 = title
                continue
            
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            # Clear placeholder text but preserve formatting
            if text_frame.paragraphs:
                text_frame.paragraphs[0].text = ""
                while len(text_frame.paragraphs) > 1:
                    elem = text_frame.paragraphs[1]._element
                    elem.getparent().remove(elem)
            if quote:
                p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                # Parse markdown formatting for quote
                parse_markdown_to_paragraph(p, quote)
                # Explicitly disable bullets and remove hanging indent
                p._element.get_or_add_pPr().buNone = None
                from lxml import etree
                buNone = etree.SubElement(p._element.get_or_add_pPr(), '{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                # Remove hanging indent - set both first line and left indent to 0
                p.level = 0
                pPr = p._element.get_or_add_pPr()
                pPr.set('indent', '0')
                pPr.set('marL', '0')
                
                if quote_citation:
                    # Strip leading bullet markers from citation
                    citation_text = quote_citation.lstrip('- ').strip()
                    # Add citation as a new paragraph
                    p = text_frame.add_paragraph()
                    # Disable bullets BEFORE setting text
                    buNone = etree.SubElement(p._element.get_or_add_pPr(), '{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                    p.text = f"\n— {citation_text}"
                    p.font.italic = True
                break


def populate_closing_slide(slide, headline, paragraph, image_path):
    """Populate closing slide with name and optional image."""
    # Find title placeholder
    title_placeholder = None
    content_placeholder = None
    
    for shape in slide.placeholders:
        ph_type = shape.placeholder_format.type
        if ph_type == 1:  # Title
            title_placeholder = shape
        elif ph_type in [2, 7]:  # Body or Object
            content_placeholder = shape
    
    if title_placeholder:
        if title_placeholder.has_text_frame:
            text_frame = title_placeholder.text_frame
            if text_frame.paragraphs:
                text_frame.paragraphs[0].text = ""
                while len(text_frame.paragraphs) > 1:
                    elem = text_frame.paragraphs[1]._element
                    elem.getparent().remove(elem)
        if headline:
            # Split headline by newlines, use first line as main text
            lines = headline.split('\n')
            title_placeholder.text = lines[0]
    
    # Add paragraph if present
    if content_placeholder:
        if content_placeholder.has_text_frame:
            content_placeholder.text_frame.clear()
        if paragraph:
            p = content_placeholder.text_frame.paragraphs[0] if content_placeholder.text_frame.paragraphs else content_placeholder.text_frame.add_paragraph()
            parse_markdown_to_paragraph(p, paragraph)
    
    # Add image if present
    if image_path:
        _pic, _ph = add_image_to_slide(slide, image_path)  # placeholder info unused for closing slide


def populate_photo_slide(slide, image_path, headline=None, hide_headline=True, paragraph=None, bullets=None, larger_image=False, prs=None):
    """Populate photo-centered slide with large image and optional headline and text.

    Args:
        slide: Slide object
        image_path: Path to image
        headline: Optional headline text
        hide_headline: If True, headline is not rendered
        paragraph: Optional paragraph text
        bullets: Optional bullets list
        larger_image: If True, attempt to scale the added image to occupy most of the slide
        prs: Optional Presentation object to read slide width/height from (used for scaling)
    """
    """Populate photo-centered slide with large image and optional headline and text.

    Args:
        slide: Slide object
        image_path: Path to image
        headline: Optional headline text
        hide_headline: If True, headline is not rendered
        paragraph: Optional paragraph text
        bullets: Optional bullets list
        larger_image: If True, attempt to scale the added image to occupy most of the slide
    """
    # Debug: print all placeholders in this slide
    print(f"  populate_photo_slide - slide placeholders:")
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            print(f"    Placeholder type {ph_type}: {shape.name}")
    
    # Add headline if not hidden
    if headline and not hide_headline:
        title_shape = None
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type == 1:  # Title placeholder
                    title_shape = shape
                    break
        
        if title_shape and title_shape.has_text_frame:
            tf = title_shape.text_frame
            tf.paragraphs[0].text = ""
            parse_markdown_to_paragraph(tf.paragraphs[0], headline)
    
    # Add paragraph/bullets to body placeholder if present
    if paragraph or bullets:
        body_shape = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.has_text_frame:
                ph_type = shape.placeholder_format.type
                if ph_type in [2, 7, 14]:  # Body, Object, or Content placeholder
                    body_shape = shape
                    break
        
        if body_shape and body_shape.has_text_frame:
            text_frame = body_shape.text_frame
            text_frame.word_wrap = True
            
            # Clear placeholder text
            if text_frame.paragraphs:
                text_frame.paragraphs[0].text = ""
                while len(text_frame.paragraphs) > 1:
                    elem = text_frame.paragraphs[1]._element
                    elem.getparent().remove(elem)
            
            # Add paragraph first if present
            if paragraph:
                p = text_frame.paragraphs[0]
                parse_markdown_to_paragraph(p, paragraph)
                # Disable bullets for paragraph text
                from lxml import etree
                pPr = p._element.get_or_add_pPr()
                # Remove any existing bullet formatting
                for buNone in pPr.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buNone'):
                    buNone.getparent().remove(buNone)
                for buAutoNum in pPr.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum'):
                    buAutoNum.getparent().remove(buAutoNum)
                for buChar in pPr.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar'):
                    buChar.getparent().remove(buChar)
                # Add buNone to explicitly disable bullets
                etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
            
            # Add bullets if present
            if bullets:
                import sys
                print(f"    populate_photo_slide bullets: {repr(bullets)}", file=sys.stderr, flush=True)
                for i, bullet in enumerate(bullets):
                    # Skip empty bullets or bracket artifacts
                    if not bullet or bullet.strip() in ['[', ']', '[]']:
                        print(f"    Skipping bullet: {repr(bullet)}", file=sys.stderr, flush=True)
                        continue
                    if i == 0 and not paragraph:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0
    
    # Add image
    pic = None
    ph_bounds = None
    if image_path:
        pic, ph_bounds = add_image_to_slide(slide, image_path, centered=True)

    # If this slide requested a larger image, scale up the picture we just added
    try:
        if pic is not None and larger_image:
            # Determine slide dimensions (prefer presentation object if passed in)
            if prs is not None:
                slide_width = prs.slide_width
                slide_height = prs.slide_height
            else:
                slide_width = 9144000
                slide_height = 6858000

            margin = Inches(1)

            # Determine bounding box to constrain the scaled image
            if ph_bounds:
                b_left = ph_bounds['left']
                b_top = ph_bounds['top']
                b_width = ph_bounds['width']
                b_height = ph_bounds['height']
            else:
                b_left = int(margin)
                b_top = int(margin)
                b_width = slide_width - (margin * 2)
                b_height = slide_height - (margin * 2)

            # Compute scaled size to fit within bounding box while preserving aspect ratio
            img_aspect = pic.width / pic.height if pic.height else 1
            target_width = b_width
            target_height = int(target_width / img_aspect)
            if target_height > b_height:
                target_height = int(b_height)
                target_width = int(target_height * img_aspect)

            pic.width = int(target_width)
            pic.height = int(target_height)

            # Center within bounding box (horizontal and vertical)
            pic.left = int(b_left + (b_width - pic.width) // 2)
            pic.top = int(b_top + (b_height - pic.height) // 2)

            print(f"    Scaled picture for larger_image: width={pic.width/914400:.2f}in height={pic.height/914400:.2f}in", file=sys.stderr)
    except Exception as e:
        print(f"    Error scaling pic for larger_image: {e}", file=sys.stderr)

    # Remove empty body/content placeholders to avoid empty text boxes showing in PowerPoint
    removed = remove_empty_body_placeholders(slide)
    if removed:
        print(f"    Removed empty placeholders: {removed}")


def populate_content_slide(slide, headline, paragraph, bullets, image_path, hide_headline, slide_class=None, larger_image=False, prs=None):
    """Populate standard content slide with headline, text, bullets, and optional image.

    Args:
        slide: Slide object
        headline: Headline text
        paragraph: Paragraph text
        bullets: Bullet list (JSON string or plain text)
        image_path: Optional path to image
        hide_headline: Whether to hide the headline
        slide_class: Optional slide class hint
        larger_image: If True, attempt to scale inserted image to occupy more slide area
        prs: Optional Presentation object used to get slide dimensions for scaling
    """
    print(f"  populate_content_slide - image_path: {repr(image_path)}, hide_headline: {hide_headline}, larger_image: {larger_image}, prs_present: {prs is not None}", file=sys.stderr, flush=True)
    
    # Determine if this is a text-only slide
    is_text_only = slide_class and 'lines' in slide_class
    
    # Find title placeholder
    if not hide_headline and headline:
        title_shape = None
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type == 1:  # Title placeholder
                    title_shape = shape
                    break
        
        if title_shape and title_shape.has_text_frame:
            # Clear and add formatted headline
            tf = title_shape.text_frame
            tf.paragraphs[0].text = ""
            parse_markdown_to_paragraph(tf.paragraphs[0], headline)
    
    # Find content placeholder for bullets or paragraph
    # Try multiple placeholder types (2=Body, 7=Object, 14=Content)
    content_shape = None
    for shape in slide.shapes:
        if shape.is_placeholder and shape.has_text_frame:
            ph_type = shape.placeholder_format.type
            if ph_type in [2, 7, 14]:  # Body, Object, or Content placeholder
                content_shape = shape
                break
    
    # If no specific content placeholder found, try any text placeholder that's not the title
    if not content_shape:
        for shape in slide.shapes:
            if shape.has_text_frame and not (shape.is_placeholder and shape.placeholder_format.type == 1):
                if hasattr(shape, 'text_frame'):
                    content_shape = shape
                    break
    
    if content_shape and content_shape.has_text_frame:
        text_frame = content_shape.text_frame
        text_frame.word_wrap = True
        
        # Clear placeholder text but preserve first paragraph's formatting
        if text_frame.paragraphs:
            # Clear text from first paragraph but keep the paragraph (preserves template formatting)
            text_frame.paragraphs[0].text = ""
            # Remove extra paragraphs
            while len(text_frame.paragraphs) > 1:
                elem = text_frame.paragraphs[1]._element
                elem.getparent().remove(elem)
        
        # For text-only slides (like template-lines), prioritize paragraph content
        if is_text_only and paragraph:
            add_formatted_text_to_frame(text_frame, paragraph)
            print(f"    Adding formatted paragraph text for text-only slide: {paragraph[:50]}...")
        # For bullet slides, handle paragraph first, then bullets
        elif not is_text_only:
            paragraph_added = False
            
            # Add paragraph text first if present
            if paragraph:
                p = text_frame.paragraphs[0]
                parse_markdown_to_paragraph(p, paragraph)
                # Explicitly disable bullets for paragraph text
                from lxml import etree
                buNone = etree.SubElement(p._element.get_or_add_pPr(), '{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                paragraph_added = True
                print(f"    Adding paragraph text: {paragraph[:50]}...")
            
            # Add bullets if present
            if bullets:
                try:
                    bullet_list = json.loads(bullets)
                    if bullet_list:  # Only add if list is not empty
                        print(f"    Adding {len(bullet_list)} bullets")
                        for i, bullet in enumerate(bullet_list):
                            if i == 0 and not paragraph_added:
                                # Use first paragraph if no paragraph text was added
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            
                            # Detect indent level from markdown-style -- prefix OR leading spaces/tabs
                            indent_level = 0
                            stripped_bullet = bullet.lstrip()
                            
                            # Check for markdown -- prefix for indentation
                            if stripped_bullet.startswith('--'):
                                # Count consecutive - characters starting from 2 (-- = level 1)
                                dash_count = len(stripped_bullet) - len(stripped_bullet.lstrip('-'))
                                indent_level = max(0, dash_count - 1)  # -- is level 1, --- is level 2, etc.
                                # Remove the - markers and any following space
                                stripped_bullet = stripped_bullet.lstrip('-').lstrip()
                            else:
                                # Fall back to space/tab detection
                                leading_whitespace = len(bullet) - len(stripped_bullet)
                                
                                # Calculate level: 2 spaces or 1 tab = 1 level
                                if '\t' in bullet[:leading_whitespace]:
                                    indent_level = bullet[:leading_whitespace].count('\t')
                                else:
                                    indent_level = leading_whitespace // 2
                            
                            # Cap at level 4 (PowerPoint supports up to 9 but 4 is reasonable)
                            indent_level = min(indent_level, 4)
                            
                            # Parse markdown formatting for bullet (use stripped text)
                            parse_markdown_to_paragraph(p, stripped_bullet)
                            p.level = indent_level
                except:
                    # If JSON parsing fails, treat as plain text
                    if not paragraph_added:
                        text_frame.text = bullets
                    else:
                        p = text_frame.add_paragraph()
                        p.text = bullets
        
        # Add paragraph if present (and no bullets handled above)
        elif paragraph:
            p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
            parse_markdown_to_paragraph(p, paragraph)
            # Explicitly disable bullets for paragraph text
            from lxml import etree
            buNone = etree.SubElement(p._element.get_or_add_pPr(), '{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
    else:
        # Debug: print placeholder info
        print(f"Warning: No content placeholder found for slide. Available placeholders:")
        for shape in slide.shapes:
            if shape.is_placeholder:
                print(f"  Type: {shape.placeholder_format.type}, Has text frame: {shape.has_text_frame}")
    
    # Add image if present
    pic = None
    ph_bounds = None
    if image_path:
        pic, ph_bounds = add_image_to_slide(slide, image_path)

    # If this slide requested a larger image, scale up the picture we just added
    try:
        if pic is not None and larger_image:
            # Determine slide dimensions (prefer presentation object if passed in)
            if prs is not None:
                slide_width = prs.slide_width
                slide_height = prs.slide_height
            else:
                slide_width = 9144000
                slide_height = 6858000

            margin = Inches(1)

            # Determine bounding box to constrain the scaled image
            if ph_bounds:
                b_left = ph_bounds['left']
                b_top = ph_bounds['top']
                b_width = ph_bounds['width']
                b_height = ph_bounds['height']
            else:
                b_left = int(margin)
                b_top = int(margin)
                b_width = slide_width - (margin * 2)
                b_height = slide_height - (margin * 2)

            # Compute scaled size to fit within bounding box while preserving aspect ratio
            img_aspect = pic.width / pic.height if pic.height else 1
            target_width = b_width
            target_height = int(target_width / img_aspect)
            if target_height > b_height:
                target_height = int(b_height)
                target_width = int(target_height * img_aspect)

            pic.width = int(target_width)
            pic.height = int(target_height)

            # Center within bounding box (horizontal and vertical)
            pic.left = int(b_left + (b_width - pic.width) // 2)
            pic.top = int(b_top + (b_height - pic.height) // 2)

            print(f"    Scaled picture for larger_image: width={pic.width/914400:.2f}in height={pic.height/914400:.2f}in", file=sys.stderr)
    except Exception as e:
        print(f"    Error scaling pic for larger_image: {e}", file=sys.stderr)

    # Remove empty body/content placeholders so empty text placeholders don't render
    removed = remove_empty_body_placeholders(slide)
    if removed:
        print(f"    Removed empty placeholders: {removed}")


def add_image_to_slide(slide, image_path, centered=False):
    """Add an image to a slide, either in a picture placeholder or positioned."""
    from PIL import Image
    
    # Clean up image path
    if image_path.startswith('/assets/'):
        image_path = 'assets' + image_path[7:]
    elif image_path.startswith('assets/'):
        pass  # Already correct
    else:
        image_path = 'assets/' + image_path
    
    # Check if file exists
    if not os.path.exists(image_path):
        print(f"Warning: Image not found: {image_path}")
        return
    
    # Get image dimensions - handle GIFs specially to preserve animation
    if image_path.lower().endswith('.gif'):
        # For GIFs, either preserve animation for small files or use a PNG fallback for large GIFs
        try:
            from PIL import Image
            gif_size = os.path.getsize(image_path)
            if gif_size > GIF_FALLBACK_THRESHOLD:
                # Convert first frame to PNG fallback to avoid embedding huge animated GIFs
                try:
                    with Image.open(image_path) as img:
                        img.seek(0)
                        fallback = tempfile.NamedTemporaryFile(delete=False, suffix='.png', prefix='pptx_gif_fallback_')
                        fallback_path = fallback.name
                        fallback.close()
                        img.convert('RGBA').save(fallback_path, 'PNG')
                        image_path = fallback_path
                        # Update dimensions from the generated PNG
                        with Image.open(fallback_path) as pf:
                            img_width, img_height = pf.size
                except Exception as e:
                    print(f"Warning: Could not create PNG fallback for GIF: {e}")
                    img_width, img_height = 800, 600
            else:
                with Image.open(image_path) as img:
                    img_width, img_height = img.size
        except Exception as e:
            print(f"Warning: Could not read GIF dimensions: {e}")
            # Use default dimensions if we can't read the file
            img_width, img_height = 800, 600
    else:
        # For other formats, use normal PIL processing
        from PIL import Image
        with Image.open(image_path) as img:
            img_width, img_height = img.size
    
    # Try to find a picture placeholder
    picture_placeholder = None
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            # Debug output
            import sys
            print(f"    Found placeholder: type={ph_type}, name={shape.name}", file=sys.stderr, flush=True)
            if ph_type == 18:  # Picture placeholder
                picture_placeholder = shape
                print(f"    -> Using this as picture placeholder!", file=sys.stderr, flush=True)
                break
    
    if picture_placeholder:
        # Get placeholder dimensions
        ph_width = picture_placeholder.width
        ph_height = picture_placeholder.height
        ph_left = picture_placeholder.left
        ph_top = picture_placeholder.top
        
        import sys
        print(f"    Picture placeholder dimensions: width={ph_width/914400:.2f}in, height={ph_height/914400:.2f}in", file=sys.stderr, flush=True)
        print(f"    Picture placeholder position: left={ph_left/914400:.2f}in, top={ph_top/914400:.2f}in", file=sys.stderr, flush=True)
        print(f"    Image dimensions: {img_width}x{img_height} pixels", file=sys.stderr, flush=True)
        
        # Use the image dimensions already determined above
        img_aspect = img_width / img_height
        ph_aspect = ph_width / ph_height
        
        # Calculate size to fit within placeholder while maintaining aspect ratio
        if img_aspect > ph_aspect:
            # Image is wider - constrain by width
            new_width = ph_width
            new_height = int(ph_width / img_aspect)
        else:
            # Image is taller - constrain by height
            new_height = ph_height
            new_width = int(ph_height * img_aspect)
        
        print(f"    Calculated image size: width={new_width/914400:.2f}in, height={new_height/914400:.2f}in", file=sys.stderr, flush=True)
        
        # Center within placeholder
        left = ph_left + (ph_width - new_width) // 2
        top = ph_top + (ph_height - new_height) // 2
        
        # Remove the placeholder and add image in its place
        sp = picture_placeholder.element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)
        
        # Send image to back so it doesn't cover text
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
        # Clean up any temporary fallback PNG we created for GIFs
        try:
            if 'fallback_path' in locals() and fallback_path:
                os.remove(fallback_path)
        except Exception:
            pass
        return pic, {'left': ph_left, 'top': ph_top, 'width': ph_width, 'height': ph_height}
    else:
        # No explicit picture placeholder; try to find a large empty content/object placeholder to place the image
        largest_ph = None
        largest_area = 0
        for shape in slide.shapes:
            if getattr(shape, 'is_placeholder', False):
                try:
                    ph_type = shape.placeholder_format.type
                except Exception:
                    continue
                # Consider Body, Object, or Content placeholders
                if ph_type in [2, 7, 14]:
                    # Only use placeholders that are empty (no text) or don't have a text frame
                    empty_text = True
                    if getattr(shape, 'has_text_frame', False):
                        txt = shape.text or ''
                        if txt.strip():
                            empty_text = False
                    if not empty_text:
                        continue
                    area = shape.width * shape.height
                    if area > largest_area:
                        largest_area = area
                        largest_ph = shape
        if largest_ph is not None and largest_area > 0:
            ph_width = largest_ph.width
            ph_height = largest_ph.height
            ph_left = largest_ph.left
            ph_top = largest_ph.top
            print(f"    Using largest empty placeholder for image: {getattr(largest_ph, 'name', None)} size={ph_width/914400:.2f}x{ph_height/914400:.2f}in", file=sys.stderr, flush=True)
            img_aspect = img_width / img_height
            ph_aspect = ph_width / ph_height
            if img_aspect > ph_aspect:
                new_width = ph_width
                new_height = int(ph_width / img_aspect)
            else:
                new_height = ph_height
                new_width = int(ph_height * img_aspect)
            left = ph_left + (ph_width - new_width) // 2
            top = ph_top + (ph_height - new_height) // 2
            # Remove placeholder and insert picture
            sp = largest_ph.element
            sp.getparent().remove(sp)
            pic = slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)
            # Send to back
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
            # Clean up any temporary fallback PNG we created for GIFs
            try:
                if 'fallback_path' in locals() and fallback_path:
                    os.remove(fallback_path)
            except Exception:
                pass
            return pic, {'left': ph_left, 'top': ph_top, 'width': ph_width, 'height': ph_height}
        else:
            # Fallback to native centered size behavior
            # Use the dimensions already determined above (img_width, img_height)
            img_width_px, img_height_px = img_width, img_height
            # Convert pixels to EMUs (English Metric Units): 1 inch = 914400 EMUs, assume 96 DPI
            dpi = 96
            img_width_emu = int(img_width_px * 914400 / dpi)
            img_height_emu = int(img_height_px * 914400 / dpi)
            
            # Get slide dimensions (standard is 10" x 7.5")
            slide_width = 9144000  # 10 inches in EMUs
            slide_height = 6858000  # 7.5 inches in EMUs
            
            # Center the image
            left = (slide_width - img_width_emu) // 2
            top = (slide_height - img_height_emu) // 2
            
            # Ensure image doesn't exceed slide bounds
            if img_width_emu > slide_width or img_height_emu > slide_height:
                # Scale down to fit
                scale = min(slide_width / img_width_emu, slide_height / img_height_emu) * 0.9
                img_width_emu = int(img_width_emu * scale)
                img_height_emu = int(img_height_emu * scale)
                left = (slide_width - img_width_emu) // 2
                top = (slide_height - img_height_emu) // 2
            
            pic = slide.shapes.add_picture(image_path, left, top, width=img_width_emu, height=img_height_emu)
            # Clean up any temporary fallback PNG we created for GIFs
            try:
                if 'fallback_path' in locals() and fallback_path:
                    os.remove(fallback_path)
            except Exception:
                pass
            return pic, None
