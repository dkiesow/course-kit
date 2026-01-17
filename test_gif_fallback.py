import json
import os
import zipfile
from pptx import Presentation
from lxml import etree
from pptx_builder import build_pptx_from_slides, GIF_FALLBACK_THRESHOLD


def test_large_gif_fallback(tmp_path):
    # Use the repo's large GIF (assets/vector.gif)
    gif_name = 'vector.gif'
    gif_path = os.path.join('assets', gif_name)
    assert os.path.exists(gif_path), "Test GIF not found in assets"

    bullets = json.dumps(["Slide with big GIF"]) 

    slides_data = [
        (None, "GIF Headline", None, bullets, None, None, gif_name, 0, False, False, False, 'bullets-image-top')
    ]

    output_file = str(tmp_path / 'gif_fallback_test.pptx')

    success = build_pptx_from_slides(
        slides_data=slides_data,
        output_path=output_file,
        template_path='templates/4734_template.potx',
        pptx_layouts_map=json.load(open('pptx_layouts.json'))
    )

    assert success
    assert os.path.exists(output_file)

    # Open with python-pptx to ensure it loads
    prs = Presentation(output_file)
    assert len(prs.slides) == 1

    # Check the slide's picture reference: if the original GIF was smaller than the threshold
    # (and may be animated), it should be preserved as a GIF; otherwise a PNG fallback should be used.
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
          'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
          'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}

    slide = prs.slides[0]
    found_pic = False
    pic_target = None
    for shape in slide.shapes:
        tag = etree.QName(shape._element.tag).localname
        if tag == 'pic':
            blip = shape._element.find('.//a:blip', namespaces=ns)
            if blip is None:
                continue
            rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            if rId and rId in slide.part.rels:
                target = slide.part.rels[rId].target_ref
                found_pic = True
                pic_target = target.lower()

    assert found_pic, "No picture found on exported slide"

    gif_size = os.path.getsize(gif_path)
    if gif_size <= GIF_FALLBACK_THRESHOLD:
        # Should be preserved as GIF (including animated GIFs)
        assert pic_target and pic_target.endswith('.gif'), f"Expected GIF to be preserved, got {pic_target}"
        # Confirm a GIF file exists in the package
        with zipfile.ZipFile(output_file, 'r') as zf:
            gifs = [n for n in zf.namelist() if n.startswith('ppt/media/') and n.lower().endswith('.gif')]
            assert gifs, "Expected a GIF to be embedded in the output PPTX"
    else:
        # Should have used a PNG fallback and no large GIFs embedded
        assert pic_target and pic_target.endswith('.png'), f"Expected PNG fallback, got {pic_target}"
        with zipfile.ZipFile(output_file, 'r') as zf:
            large_gifs = [(n, zf.getinfo(n).file_size) for n in zf.namelist() if n.startswith('ppt/media/') and n.lower().endswith('.gif') and zf.getinfo(n).file_size > GIF_FALLBACK_THRESHOLD]
            assert not large_gifs, f"Large GIFs were embedded in PPTX: {large_gifs}"
