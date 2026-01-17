import json
import os
from pptx import Presentation
from pptx_builder import build_pptx_from_slides


def test_fullscreen_uses_top_bullets_layout(tmp_path):
    # Prepare slides_data: one slide with bullets, image, and fullscreen=True
    image_name = 'milkshake.png'
    assert os.path.exists(os.path.join('assets', image_name)), "Test image not found in assets"

    bullets = json.dumps(["First bullet", "Second bullet"])

    slides_data = [
        (None, "Headline", None, bullets, None, None, image_name, 0, False, False, True, 'bullets-image-top')
    ]

    output_file = str(tmp_path / "test_fullscreen.pptx")

    success = build_pptx_from_slides(
        slides_data=slides_data,
        output_path=output_file,
        template_path='templates/4734_template.potx',
        pptx_layouts_map=json.load(open('pptx_layouts.json'))
    )

    assert success
    assert os.path.exists(output_file), "PPTX output not created"

    prs = Presentation(output_file)
    # There should be one slide
    assert len(prs.slides) == 1

    layout_name = prs.slides[0].slide_layout.name
    assert layout_name == "White_Full_Photo_Headline", f"Unexpected layout: {layout_name}"

    # Verify headline is present in the title placeholder (not hidden by the image)
    title_text = ''
    for shape in prs.slides[0].shapes:
        if shape.is_placeholder and shape.placeholder_format.type == 1 and shape.has_text_frame:
            title_text = shape.text.strip()
            break
    assert 'Headline' in title_text, f"Headline not found in title placeholder: {title_text}"

    # Verify picture element appears before title in the slide XML (i.e., picture is behind title)
    from lxml import etree
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    spTree = prs.slides[0]._element.find('.//p:cSld/p:spTree', namespaces=ns)
    children = list(spTree)

    pic_index = None
    title_index = None
    for i, ch in enumerate(children):
        tag = etree.QName(ch.tag).localname
        if tag == 'pic' and pic_index is None:
            pic_index = i
        if tag == 'sp' and title_index is None:
            ph = ch.find('.//p:ph', namespaces=ns)
            if ph is not None and ph.get('type') == 'title':
                title_index = i
    assert pic_index is not None and title_index is not None, 'Could not find pic or title elements in slide XML'
    assert pic_index < title_index, f'Picture element ({pic_index}) is not behind title ({title_index})'


def test_fullscreen_does_not_force_white_for_gold(tmp_path):
    # Gold templates should not be forced to white layouts even if fullscreen is checked
    image_name = 'milkshake.png'
    bullets = json.dumps(["First bullet", "Second bullet"])

    slides_data = [
        (None, "Headline", None, bullets, None, None, image_name, 0, False, False, True, 'gold-bullets-image-top')
    ]

    output_file = str(tmp_path / "test_fullscreen_gold.pptx")

    success = build_pptx_from_slides(
        slides_data=slides_data,
        output_path=output_file,
        template_path='templates/4734_template.potx',
        pptx_layouts_map=json.load(open('pptx_layouts.json'))
    )

    assert success
    prs = Presentation(output_file)
    layout_name = prs.slides[0].slide_layout.name
    # Should be a Gold layout
    assert layout_name.startswith('Gold_'), f"Gold template unexpectedly used white layout: {layout_name}"


def test_full_photo_headline_template_mapping(tmp_path):
    # The 'full-photo-headline' template should map to White_Top_Bullets_Photo
    image_name = 'milkshake.png'
    bullets = json.dumps(["First bullet", "Second bullet"])

    slides_data = [
        (None, "Headline", None, bullets, None, None, image_name, 0, False, False, False, 'full-photo-headline')
    ]

    output_file = str(tmp_path / "test_full_photo_headline.pptx")

    success = build_pptx_from_slides(
        slides_data=slides_data,
        output_path=output_file,
        template_path='templates/4734_template.potx',
        pptx_layouts_map=json.load(open('pptx_layouts.json'))
    )

    assert success
    prs = Presentation(output_file)
    layout_name = prs.slides[0].slide_layout.name
    assert layout_name == "White_Full_Photo_Headline", f"Unexpected layout for full-photo-headline: {layout_name}"

    # Verify headline is present
    title_text = ''
    for shape in prs.slides[0].shapes:
        if shape.is_placeholder and shape.placeholder_format.type == 1 and shape.has_text_frame:
            title_text = shape.text.strip()
            break
    assert 'Headline' in title_text, f"Headline not found in title placeholder: {title_text}"

    # Also assert picture is behind title in XML ordering
    from lxml import etree
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    spTree = prs.slides[0]._element.find('.//p:cSld/p:spTree', namespaces=ns)
    children = list(spTree)

    pic_index = None
    title_index = None
    for i, ch in enumerate(children):
        tag = etree.QName(ch.tag).localname
        if tag == 'pic' and pic_index is None:
            pic_index = i
        if tag == 'sp' and title_index is None:
            ph = ch.find('.//p:ph', namespaces=ns)
            if ph is not None and ph.get('type') == 'title':
                title_index = i
    assert pic_index is not None and title_index is not None
    assert pic_index < title_index, 'Picture is not behind the title (ordering mismatch)'
