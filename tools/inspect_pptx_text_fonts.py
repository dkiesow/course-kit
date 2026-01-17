#!/usr/bin/env python3
"""Inspect text and fonts in a PPTX to help match PowerPoint's docProps counting."""
import sys
from pptx import Presentation


def count_words_in_text(s):
    if s is None:
        return 0
    return len([w for w in s.strip().split() if w.strip()])


def inspect(pptx_path):
    prs = Presentation(pptx_path)
    total_paragraphs = 0
    total_words = 0
    per_slide = []
    fonts = set()
    for i, s in enumerate(prs.slides, start=1):
        slide_words = 0
        slide_paras = 0
        # shapes text and tables
        for shape in s.shapes:
            try:
                if getattr(shape, 'has_text_frame', False):
                    for p in shape.text_frame.paragraphs:
                        if p.text and p.text.strip():
                            slide_paras += 1
                            w = count_words_in_text(p.text)
                            slide_words += w
                            # runs fonts
                            for run in p.runs:
                                try:
                                    fn = getattr(run.font, 'name', None)
                                    if fn:
                                        fonts.add(fn)
                                except Exception:
                                    pass
                if getattr(shape, 'has_table', False):
                    tbl = shape.table
                    for row in tbl.rows:
                        for cell in row.cells:
                            if cell.text and cell.text.strip():
                                slide_paras += 1
                                w = count_words_in_text(cell.text)
                                slide_words += w
                                # run fonts in cell
                                try:
                                    for p in cell.text_frame.paragraphs:
                                        for run in p.runs:
                                            fn = getattr(run.font, 'name', None)
                                            if fn:
                                                fonts.add(fn)
                                except Exception:
                                    pass
            except Exception:
                pass
        # note slides
        try:
            ns = s.notes_slide
            if ns:
                for shape in ns.shapes:
                    try:
                        if getattr(shape, 'has_text_frame', False):
                            for p in shape.text_frame.paragraphs:
                                if p.text and p.text.strip():
                                    slide_paras += 1
                                    w = count_words_in_text(p.text)
                                    slide_words += w
                                    for run in p.runs:
                                        fn = getattr(run.font, 'name', None)
                                        if fn:
                                            fonts.add(fn)
                    except Exception:
                        pass
        except Exception:
            pass

        total_paragraphs += slide_paras
        total_words += slide_words
        per_slide.append((i, slide_paras, slide_words))

    print('slides:', len(prs.slides))
    print('paragraphs:', total_paragraphs)
    print('words:', total_words)
    print('\nTop slides by words:')
    for s in sorted(per_slide, key=lambda x: x[2], reverse=True)[:10]:
        print(f' slide {s[0]:>2}: paragraphs={s[1]:>3} words={s[2]:>4}')
    print('\nFonts found (by run.font.name):')
    for f in sorted(fonts):
        print(' -', f)


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('Usage: inspect_pptx_text_fonts.py <pptx>')
        sys.exit(2)
    inspect(sys.argv[1])
