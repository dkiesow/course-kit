import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx_builder import normalize_pptx


def test_normalize_populates_docprops(tmp_path):
    prs = Presentation()
    # Title slide
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    try:
        if s0.shapes.title and s0.shapes.title.has_text_frame:
            s0.shapes.title.text = 'Test Title'
    except Exception:
        pass

    # Content slide with two paragraphs
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    try:
        body = s1.placeholders[1]
        body.text = 'First paragraph\nSecond paragraph'
    except Exception:
        # Fallback: add a textbox with text
        tx = s1.shapes.add_textbox(91440, 91440, 4000000, 1000000)
        tx.text = 'First paragraph\nSecond paragraph'

    p = str(tmp_path / 'test_pres.pptx')
    prs.save(p)

    # Run normalization which should update docProps/app.xml
    normalize_pptx(p)

    with zipfile.ZipFile(p) as zf:
        app_xml = zf.read('docProps/app.xml').decode('utf-8')
        tree = ET.fromstring(app_xml)
        ns = 'http://schemas.openxmlformats.org/officeDocument/2006/extended-properties'
        words = tree.find(f'{{{ns}}}Words')
        paras = tree.find(f'{{{ns}}}Paragraphs')
        slides = tree.find(f'{{{ns}}}Slides')
        assert words is not None and int(words.text) > 0
        assert paras is not None and int(paras.text) > 0
        assert slides is not None and int(slides.text) >= 2

        # TitlesOfParts should list at least one known title
        tops = tree.find(f'{{{ns}}}TitlesOfParts')
        assert tops is not None
        tops_text = ''.join([elt.text or '' for elt in tops.iter()])
        assert 'Test Title' in tops_text or 'PowerPoint Presentation' in tops_text