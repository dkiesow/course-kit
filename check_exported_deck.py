#!/usr/bin/env python3
"""
Test the exported PPTX to verify font inheritance is working correctly
"""

from pptx import Presentation

def test_exported_pptx():
    pptx_file = 'output/Week_Week One Wednesday_1-21-26.pptx'
    
    print("=== TESTING EXPORTED PPTX ===")
    
    # Load the presentation
    prs = Presentation(pptx_file)
    
    print(f"Number of slides: {len(prs.slides)}")
    
    # Check a few slides for layout types
    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        print(f"Slide {i+1}: {layout_name}")
        
        # Look for Gold vs White layouts
        if 'Gold' in layout_name:
            print(f"  *** GOLD LAYOUT FOUND ***")
        elif 'White' in layout_name:
            print(f"  *** WHITE LAYOUT FOUND ***")
        
        # Stop after checking first 10 slides
        if i >= 9:
            print("... (truncated)")
            break
    
    print(f"\nExported PPTX file: {pptx_file}")
    print("Open this file to verify:")
    print("- Gold layouts have gold backgrounds") 
    print("- Gold layouts use Helvetica Neue fonts (not Aptos)")
    print("- White layouts work as expected")

if __name__ == "__main__":
    test_exported_pptx()