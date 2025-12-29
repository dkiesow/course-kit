from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
import zipfile
import tempfile
import os

def convert_potx_to_pptx(potx_path):
    """Convert POTX to PPTX by modifying content type"""
    with tempfile.NamedTemporaryFile(mode='w+b', suffix='.pptx', delete=False) as tmp:
        tmp_path = tmp.name
    
    with zipfile.ZipFile(potx_path, 'r') as zip_in:
        with zipfile.ZipFile(tmp_path, 'w') as zip_out:
            for item in zip_in.infolist():
                data = zip_in.read(item.filename)
                if item.filename == '[Content_Types].xml':
                    data = data.replace(
                        b'application/vnd.openxmlformats-officedocument.presentationml.template.main',
                        b'application/vnd.openxmlformats-officedocument.presentationml.presentation.main'
                    )
                zip_out.writestr(item, data)
    
    return tmp_path

# Load template
template_path = '/Users/kiesowd/slides-marp/4734_template.potx'
tmp_pptx = convert_potx_to_pptx(template_path)

try:
    prs = Presentation(tmp_pptx)
    
    print("Master slide placeholders:\n")
    
    for master_idx, master in enumerate(prs.slide_masters):
        print(f"\nMaster {master_idx}:")
        print(f"  Total placeholders: {len(master.placeholders)}")
        
        for ph_idx, ph in enumerate(master.placeholders):
            try:
                ph_type = ph.placeholder_format.type
                type_name = PP_PLACEHOLDER(ph_type).name if ph_type in [e.value for e in PP_PLACEHOLDER] else f"UNKNOWN_{ph_type}"
            except:
                type_name = "CANNOT_DETERMINE"
            
            has_text = hasattr(ph, 'text_frame')
            print(f"    idx={ph.placeholder_format.idx}, type={type_name} ({ph_type}), has_text_frame={has_text}")
        
        # Now check the layouts under this master
        print(f"  Layouts under this master:")
        for layout_idx, layout in enumerate(master.slide_layouts):
            print(f"    Layout: {layout.name}")
            print(f"      Layout placeholders: {len(layout.placeholders)}")
            for ph in layout.placeholders:
                try:
                    ph_type = ph.placeholder_format.type
                    type_name = PP_PLACEHOLDER(ph_type).name if ph_type in [e.value for e in PP_PLACEHOLDER] else f"UNKNOWN_{ph_type}"
                except:
                    type_name = "CANNOT_DETERMINE"
                print(f"        idx={ph.placeholder_format.idx}, type={type_name} ({ph_type})")
            
            # Create a slide from this layout and check
            slide = prs.slides.add_slide(layout)
            print(f"      Slide created from layout has {len(slide.placeholders)} placeholders:")
            for ph in slide.placeholders:
                try:
                    ph_type = ph.placeholder_format.type
                    type_name = PP_PLACEHOLDER(ph_type).name if ph_type in [e.value for e in PP_PLACEHOLDER] else f"UNKNOWN_{ph_type}"
                except:
                    type_name = "CANNOT_DETERMINE"
                print(f"        idx={ph.placeholder_format.idx}, type={type_name} ({ph_type})")

finally:
    os.unlink(tmp_pptx)
