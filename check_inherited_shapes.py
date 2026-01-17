from pptx import Presentation
import shutil
import zipfile

# Patch POTX to work with python-pptx
template_path = 'templates/4734_template.potx'
temp_path = 'temp_template.pptx'

shutil.copy(template_path, temp_path)

# Modify content type
with zipfile.ZipFile(temp_path, 'r') as zip_read:
    content_types = zip_read.read('[Content_Types].xml').decode('utf-8')
    content_types = content_types.replace(
        'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
    )
    
    with zipfile.ZipFile(temp_path + '.new', 'w', zipfile.ZIP_DEFLATED) as zip_write:
        for item in zip_read.infolist():
            data = zip_read.read(item.filename)
            if item.filename == '[Content_Types].xml':
                zip_write.writestr(item, content_types.encode('utf-8'))
            else:
                zip_write.writestr(item, data)

shutil.move(temp_path + '.new', temp_path)

prs = Presentation(temp_path)

# Get White theme master and White_Bullets layout
master = prs.slide_masters[0]
layout = None
for l in master.slide_layouts:
    if l.name == "White_Bullets":
        layout = l
        break

print("=== Checking if we can access master placeholders from layout ===")
print(f"\nLayout: {layout.name}")
print(f"Layout has {len(layout.placeholders)} placeholders")
print(f"Master has {len(master.placeholders)} placeholders")

print("\n=== Can we iterate master placeholders? ===")
for placeholder in master.placeholders:
    print(f"Master placeholder idx={placeholder.placeholder_format.idx}: type={placeholder.placeholder_format.type}, name={placeholder.name}")

print("\n=== Create a slide and check ALL shapes (not just placeholders) ===")
slide = prs.slides.add_slide(layout)
print(f"Slide has {len(slide.shapes)} shapes")

# Check each shape
for i, shape in enumerate(slide.shapes):
    print(f"\nShape {i}:")
    print(f"  Name: {shape.name}")
    print(f"  Type: {shape.shape_type}")
    print(f"  Has text frame: {shape.has_text_frame}")
    if shape.has_text_frame:
        print(f"  Text frame paragraphs: {len(shape.text_frame.paragraphs)}")
    
    # Check if this shape matches the master's body placeholder
    if shape.has_text_frame and shape.is_placeholder:
        print(f"  Is placeholder: True")
        print(f"  Placeholder idx: {shape.placeholder_format.idx}")
        print(f"  Placeholder type: {shape.placeholder_format.type}")
        
# Check if we can manually add a text box in the body area
print("\n=== Checking if we can manually add text to slide ===")
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

# Try adding a text box where the body content should go
left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(4)

textbox = slide.shapes.add_textbox(left, top, width, height)
textbox.text = "This is manually added body text"

print(f"Added textbox successfully")
print(f"Slide now has {len(slide.shapes)} shapes")

# Save to check
prs.save('test_manual_text.pptx')
print("\nSaved test_manual_text.pptx to verify")
