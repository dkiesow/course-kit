#!/usr/bin/env python3
"""
Debug font availability and substitution during PPTX creation
"""

from pptx import Presentation
import shutil
import zipfile

# Test font availability and what actually gets applied
def debug_font_substitution():
    template_path = 'templates/4734_template.potx'
    temp_path = 'temp_template.pptx'
    
    # Patch POTX to PPTX
    shutil.copy(template_path, temp_path)
    
    with zipfile.ZipFile(temp_path, 'r') as zip_read:
        content_types = zip_read.read('[Content_Types].xml').decode('utf-8')
        content_types = content_types.replace(
            'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
        )
        
        with zipfile.ZipFile(temp_path + '.new', 'w', zipfile.ZIP_DEFLATED) as zip_write:
            for item in zip_read.infolist():
                data = zip_read.read(item.filename)
                if item.filename == '[Content_Types].xml':
                    zip_write.writestr(item, content_types.encode('utf-8'))
                else:
                    zip_write.writestr(item, data)
    
    shutil.move(temp_path + '.new', temp_path)
    
    # Load and examine presentation
    prs = Presentation(temp_path)
    
    print("=== FONT DEBUG ANALYSIS ===")
    
    # Check available fonts in the system (what python-pptx can see)
    print("\nTemplate font analysis:")
    
    layouts = {}
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            layouts[layout.name] = layout
    
    # Create test slides for both White and Gold and compare
    slides_created = []
    
    for layout_name in ['White_Bullets', 'Gold_Bullets']:
        if layout_name in layouts:
            slide = prs.slides.add_slide(layouts[layout_name])
            title = slide.shapes.title
            title.text = f"{layout_name} Test Title"
            slides_created.append(layout_name)
            
            # Add some bullet text
            content = None
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape != title:
                    content = shape
                    break
            
            if content:
                content.text = f"{layout_name} bullet point"
                
                print(f"\n{layout_name} slide created:")
                print(f"Title shape:")
                if hasattr(title, 'text_frame'):
                    for paragraph in title.text_frame.paragraphs:
                        for run in paragraph.runs:
                            font = run.font
                            print(f"  Font name: {font.name}")
                            print(f"  Font size: {font.size}")
                            print(f"  Font bold: {font.bold}")
                
                print(f"Content shape:")
                if hasattr(content, 'text_frame'):
                    for paragraph in content.text_frame.paragraphs:
                        for run in paragraph.runs:
                            font = run.font
                            print(f"  Font name: {font.name}")
                            print(f"  Font size: {font.size}")
                            print(f"  Font bold: {font.bold}")
    
    # Save and examine the generated file
    output_path = 'debug_fonts_test.pptx'
    prs.save(output_path)
    
    print(f"\nSaved test file: {output_path}")
    print("Check this file to see if fonts are Helvetica Neue or Aptos")
    
    # Also examine the XML of the created slide
    print("\n=== EXAMINING GENERATED SLIDE XML ===")
    with zipfile.ZipFile(output_path, 'r') as zip_read:
        # Look at slide1.xml
        try:
            slide_xml = zip_read.read('ppt/slides/slide1.xml').decode('utf-8')
            
            # Look for font references in the slide
            if 'Helvetica' in slide_xml:
                print("✓ Found Helvetica in slide XML")
            else:
                print("✗ No Helvetica found in slide XML")
                
            if 'Aptos' in slide_xml:
                print("✗ Found Aptos in slide XML")
            else:
                print("✓ No Aptos found in slide XML")
                
            # Look for theme font references
            if '+mn-lt' in slide_xml:
                print("Found +mn-lt theme font reference")
            if '+mj-lt' in slide_xml:
                print("Found +mj-lt theme font reference")
                
        except Exception as e:
            print(f"Error reading slide XML: {e}")

if __name__ == "__main__":
    debug_font_substitution()